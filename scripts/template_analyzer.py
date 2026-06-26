"""
PPTX 模板分析器 (Template Analyzer)
=====================================
读取现有 .pptx 文件的母版信息，提取设计令牌供 SVG 生成使用。

依赖：
  pip install python-pptx

用法：
  python template-analyzer.py <template.pptx>

输出：
  - 控制台打印设计令牌摘要
  - 在同目录生成 extracted-tokens.json（机器可读）
"""

import sys
import json
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn


def emu_to_px(emu, dpi=96):
    """EMU -> pixels (at 96 DPI)"""
    return round(emu / 914400 * dpi, 1)


def emu_to_inches(emu):
    """EMU -> inches"""
    return round(emu / 914400, 3)


def rgb_to_hex(r, g, b):
    return f"#{r:02X}{g:02X}{b:02X}"


def extract_theme_colors(prs):
    """提取主题配色方案。"""
    colors = {}
    try:
        theme = prs.slide_masters[0].element.find(
            './/' + qn('a:clrScheme'))
        if theme is not None:
            color_names = [
                'dk1', 'lt1', 'dk2', 'lt2',
                'accent1', 'accent2', 'accent3', 'accent4',
                'accent5', 'accent6', 'hlink', 'folHlink'
            ]
            display_names = {
                'dk1': '深色1 (文本)',
                'lt1': '浅色1 (背景)',
                'dk2': '深色2',
                'lt2': '浅色2',
                'accent1': '强调色1',
                'accent2': '强调色2',
                'accent3': '强调色3',
                'accent4': '强调色4',
                'accent5': '强调色5',
                'accent6': '强调色6',
                'hlink': '链接色',
                'folHlink': '已访问链接色'
            }
            for name in color_names:
                elem = theme.find(qn(f'a:{name}'))
                if elem is not None:
                    # Try srgbClr first
                    srgb = elem.find(qn('a:srgbClr'))
                    if srgb is not None:
                        hex_val = '#' + srgb.get('val', '000000')
                        colors[name] = {
                            'hex': hex_val,
                            'display': display_names.get(name, name)
                        }
                    else:
                        # Try sysClr
                        sys_clr = elem.find(qn('a:sysClr'))
                        if sys_clr is not None:
                            last_clr = sys_clr.get('lastClr', '000000')
                            colors[name] = {
                                'hex': '#' + last_clr,
                                'display': display_names.get(name, name),
                                'system': sys_clr.get('val', '')
                            }
    except Exception as e:
        print(f"  [WARN] 提取主题配色失败: {e}")
    return colors


def extract_theme_fonts(prs):
    """提取主题字体。"""
    fonts = {}
    try:
        font_scheme = prs.slide_masters[0].element.find(
            './/' + qn('a:fontScheme'))
        if font_scheme is not None:
            fonts['scheme_name'] = font_scheme.get('name', 'Unknown')

            major = font_scheme.find(qn('a:majorFont'))
            if major is not None:
                latin = major.find(qn('a:latin'))
                ea = major.find(qn('a:ea'))
                fonts['heading_latin'] = latin.get('typeface', '') if latin is not None else ''
                fonts['heading_ea'] = ea.get('typeface', '') if ea is not None else ''

            minor = font_scheme.find(qn('a:minorFont'))
            if minor is not None:
                latin = minor.find(qn('a:latin'))
                ea = minor.find(qn('a:ea'))
                fonts['body_latin'] = latin.get('typeface', '') if latin is not None else ''
                fonts['body_ea'] = ea.get('typeface', '') if ea is not None else ''
    except Exception as e:
        print(f"  [WARN] 提取主题字体失败: {e}")
    return fonts


def extract_layouts(prs):
    """提取所有可用版式。"""
    layouts = []
    try:
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                placeholders = []
                for ph in layout.placeholders:
                    placeholders.append({
                        'idx': ph.placeholder_format.idx,
                        'type': str(ph.placeholder_format.type),
                        'name': ph.name,
                        'left': emu_to_inches(ph.left) if ph.left else None,
                        'top': emu_to_inches(ph.top) if ph.top else None,
                        'width': emu_to_inches(ph.width) if ph.width else None,
                        'height': emu_to_inches(ph.height) if ph.height else None,
                    })
                layouts.append({
                    'name': layout.name,
                    'placeholders': placeholders
                })
    except Exception as e:
        print(f"  [WARN] 提取版式失败: {e}")
    return layouts


def extract_slide_size(prs):
    """提取幻灯片尺寸。"""
    return {
        'width_inches': emu_to_inches(prs.slide_width),
        'height_inches': emu_to_inches(prs.slide_height),
        'width_emu': prs.slide_width,
        'height_emu': prs.slide_height,
        'aspect_ratio': f"{round(prs.slide_width / prs.slide_height, 2)}:1"
    }


def analyze_template(filepath):
    """分析 .pptx 模板并提取设计令牌。"""
    print(f"\n{'='*60}")
    print(f"  PPTX 模板分析: {filepath}")
    print(f"{'='*60}\n")

    prs = Presentation(filepath)

    # 1. 幻灯片尺寸
    size = extract_slide_size(prs)
    print(f"📐 幻灯片尺寸")
    print(f"   {size['width_inches']}\" × {size['height_inches']}\"")
    print(f"   宽高比: {size['aspect_ratio']}")
    print()

    # 2. 主题配色
    colors = extract_theme_colors(prs)
    print(f"🎨 主题配色方案")
    for key, val in colors.items():
        print(f"   {val['display']:20s} {val['hex']}")
    print()

    # 3. 主题字体
    fonts = extract_theme_fonts(prs)
    print(f"🔤 主题字体")
    print(f"   方案名称:  {fonts.get('scheme_name', 'N/A')}")
    print(f"   标题(西文): {fonts.get('heading_latin', 'N/A')}")
    print(f"   标题(中文): {fonts.get('heading_ea', 'N/A')}")
    print(f"   正文(西文): {fonts.get('body_latin', 'N/A')}")
    print(f"   正文(中文): {fonts.get('body_ea', 'N/A')}")
    print()

    # 4. 可用版式
    layouts = extract_layouts(prs)
    print(f"📋 可用版式 ({len(layouts)} 个)")
    for i, layout in enumerate(layouts):
        ph_summary = ', '.join([p['name'] for p in layout['placeholders']])
        print(f"   [{i}] {layout['name']}")
        if ph_summary:
            print(f"       占位符: {ph_summary}")
    print()

    # 5. 现有幻灯片数
    print(f"📄 现有幻灯片: {len(prs.slides)} 页")
    print()

    # 输出 JSON
    tokens = {
        'slide_size': size,
        'colors': colors,
        'fonts': fonts,
        'layouts': layouts,
        'existing_slides': len(prs.slides)
    }

    output_path = filepath.rsplit('.', 1)[0] + '_tokens.json'
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(tokens, f, ensure_ascii=False, indent=2)
    print(f"💾 设计令牌已保存: {output_path}")

    # 打印可直接用于 SVG 的配色建议
    print(f"\n{'─'*60}")
    print(f"📌 建议的 SVG 设计令牌映射")
    print(f"{'─'*60}")
    if 'lt1' in colors:
        print(f'   fill-bg      = "{colors["lt1"]["hex"]}"  (背景)')
    if 'accent1' in colors:
        print(f'   fill-accent   = "{colors["accent1"]["hex"]}"  (强调色/焦点)')
    if 'dk1' in colors:
        print(f'   fill-title    = "{colors["dk1"]["hex"]}"  (标题)')
    if 'dk2' in colors:
        print(f'   fill-body     = "{colors["dk2"]["hex"]}"  (正文)')
    if 'lt2' in colors:
        print(f'   fill-card     = "{colors["lt2"]["hex"]}"  (卡片背景)')
    heading_font = fonts.get('heading_ea') or fonts.get('heading_latin', 'Microsoft YaHei')
    body_font = fonts.get('body_ea') or fonts.get('body_latin', 'Microsoft YaHei')
    print(f'   font-heading  = "{heading_font}"')
    print(f'   font-body     = "{body_font}"')
    print()

    return tokens


def main():
    if len(sys.argv) < 2:
        print("用法: python template-analyzer.py <template.pptx>")
        print("\n分析 .pptx 文件的母版，提取设计令牌供 SVG 生成使用。")
        sys.exit(1)

    filepath = sys.argv[1]
    try:
        analyze_template(filepath)
    except Exception as e:
        print(f"\n❌ 分析失败: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
