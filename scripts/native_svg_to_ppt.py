"""
Smart SVG -> PPTX Native Converter (v5.0)
==========================================
v5.0 改进 (在 v4.3 基础上)：
- 🔴 新增：<path> 元素支持（M/L/Z 直线 + C/Q/A 曲线折线近似）
- 🔴 新增：<linearGradient> / <radialGradient> 解析，回退为首 stop-color 纯色
- 🔴 修复：fill="url(#id)" 自动查 gradient_map 解析为纯色
- 🔴 新增：fill-opacity / stroke-opacity 属性支持
- 🔴 新增：<g transform="scale(sx,sy)"> 支持
- 🟡 新增：<ellipse> 元素支持

v4.3 原有：
- clip-path="inset(... round R1 R2 R3 R4)" 提取圆角
- polygon/freeform fill XML 注入（p:spPr 命名空间查找）
- letter-spacing 属性支持 → PPT a:rPr spc
- opacity 属性支持 → PPT 形状透明度

v4.2 原有：
- polygon/freeform 使用 XML 注入 solidFill
- <defs><clipPath> rx 解析
- connector 阴影移除
- 全画布背景 rect → slide background
- <g> 继承 stroke/stroke-width
- 演讲者备注写入
"""

import os
import sys
import glob
import json
import re
import math
import argparse
import xml.etree.ElementTree as ET
from urllib.parse import urlparse, unquote

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 画布参数 ──
SVG_W = 1920
SVG_H = 1080
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SCALE = SLIDE_W_IN / SVG_W
FONT_SCALE = SCALE * 72

SLIDE_W = Inches(SLIDE_W_IN)
SLIDE_H = Inches(SLIDE_H_IN)


# ═══════════════════════════════════════
# 画布/比例动态配置
# ═══════════════════════════════════════

def set_canvas(svg_w, svg_h, slide_w_in=13.333, slide_h_in=7.5):
    """动态设置 SVG 画布尺寸与 PPT 页尺寸，并同步 SCALE/FONT_SCALE。"""
    global SVG_W, SVG_H, SLIDE_W_IN, SLIDE_H_IN, SCALE, FONT_SCALE, SLIDE_W, SLIDE_H
    SVG_W = float(svg_w)
    SVG_H = float(svg_h)
    SLIDE_W_IN = float(slide_w_in)
    SLIDE_H_IN = float(slide_h_in)
    SCALE = SLIDE_W_IN / SVG_W
    FONT_SCALE = SCALE * 72
    SLIDE_W = Inches(SLIDE_W_IN)
    SLIDE_H = Inches(SLIDE_H_IN)


def parse_svg_canvas_size(svg_file):
    """从 SVG 的 viewBox 或 width/height 推断画布尺寸。失败则回退当前全局 SVG_W/H。"""
    try:
        tree = ET.parse(svg_file)
        root = tree.getroot()
        vb = root.attrib.get('viewBox') or root.attrib.get('viewbox')
        if vb:
            parts = [float(x) for x in re.split(r'[\s,]+', vb.strip()) if x]
            if len(parts) == 4:
                return float(parts[2]), float(parts[3])
        w = root.attrib.get('width')
        h = root.attrib.get('height')
        if w and h:
            wv = float(re.sub(r'[^\d.]+', '', str(w)))
            hv = float(re.sub(r'[^\d.]+', '', str(h)))
            if wv > 0 and hv > 0:
                return wv, hv
    except:
        pass
    return float(SVG_W), float(SVG_H)


# ═══════════════════════════════════════
# 渐变解析 (v5.0)
# ═══════════════════════════════════════

def parse_gradients(root):
    """预扫描 <defs> 中所有 linearGradient / radialGradient，
    构建 gradient_map: {id → first_stop_color_hex}。
    PPT 不支持 SVG 渐变，用第 1 个 stop-color 作为纯色回退。
    """
    gradient_map = {}
    for node in root.iter():
        tag = node.tag.split('}')[-1] if '}' in node.tag else node.tag
        if tag in ('linearGradient', 'radialGradient'):
            gid = node.attrib.get('id')
            if gid:
                stops = []
                for child in node:
                    ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if ctag == 'stop':
                        color = child.attrib.get('stop-color', '')
                        offset = child.attrib.get('offset', '0%')
                        stops.append((offset, color))
                if stops:
                    # 使用第一个 stop-color 作为回退色
                    gradient_map[gid] = stops[0][1]
    return gradient_map


# ═══════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════

# 全局 gradient_map，在 main() 中每页设置
_gradient_map = {}
_current_svg_dir = None
_conversion_errors = []
_conversion_warnings = []
_strict_missing_images = False


def warn(msg):
    """Record a non-blocking conversion warning and print it for CLI users."""
    _conversion_warnings.append(msg)
    print(f'  WARN: {msg}')


def error(msg):
    """Record a blocking conversion error and print it for CLI users."""
    _conversion_errors.append(msg)
    print(f'  ERROR: {msg}')

def parse_color(c):
    """解析 SVG 颜色值为 RGBColor，不支持的返回 None。
    v5.0: 支持 url(#id) 引用 gradient_map 回退。
    """
    if not c or c == 'none':
        return None
    c = c.strip()
    # v5.0: url(#id) 引用渐变 → 回退首色
    if c.startswith('url('):
        m = re.match(r'url\(#(.+?)\)', c)
        if m:
            gid = m.group(1)
            fallback = _gradient_map.get(gid)
            if fallback:
                return parse_color(fallback)
        return None
    if c.startswith('rgb('):
        try:
            inner = c[4:].rstrip(')')
            parts = [int(x.strip()) for x in inner.split(',')]
            return RGBColor(parts[0], parts[1], parts[2])
        except:
            return None
    c = c.lstrip('#')
    if len(c) == 3:
        c = c[0]*2 + c[1]*2 + c[2]*2
    if len(c) == 6:
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    return None


def resolve_image_href(href):
    """将 SVG <image> 的 href 解析为本地文件路径。仅支持本地文件（file:// 或相对路径）。"""
    global _current_svg_dir
    if not href:
        return None
    href = href.strip()
    if href.startswith('data:'):
        return None

    # file:// URL
    if href.startswith('file:'):
        try:
            parsed = urlparse(href)
            p = unquote(parsed.path or '')
            # Windows file URL 常见形式：/c:/path...
            if re.match(r'^/[A-Za-z]:', p):
                p = p[1:]
            p = p.replace('/', os.sep)
            return p
        except:
            return None

    # 绝对路径
    if os.path.isabs(href):
        return href

    # 相对路径
    if _current_svg_dir:
        return os.path.join(_current_svg_dir, href)
    return href


def remove_shadow(shape):
    """移除 python-pptx 形状/连接器的默认阴影效果。"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = sp.find(qn('a:spPr'))
    if spPr is None:
        for child in sp:
            if child.tag.endswith('spPr') or child.tag.endswith('cxnSpPr'):
                spPr = child
                break
    if spPr is not None:
        for eff in spPr.findall(qn('a:effectLst')):
            spPr.remove(eff)
        etree.SubElement(spPr, qn('a:effectLst'))


def _find_spPr(shape):
    """在 shape._element 中查找 spPr 节点。"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        for child in sp:
            if child.tag.endswith('spPr'):
                spPr = child
                break
    return spPr


def _apply_alpha(spPr, alpha_val):
    """向 spPr 中的 solidFill > srgbClr 注入 alpha 透明度。
    alpha_val: 0.0 (全透明) ~ 1.0 (不透明)
    """
    if spPr is None or alpha_val is None or alpha_val >= 1.0:
        return
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            # 移除已有 alpha
            for old_a in clr.findall(qn('a:alpha')):
                clr.remove(old_a)
            alpha_el = etree.SubElement(clr, qn('a:alpha'))
            alpha_el.set('val', str(int(alpha_val * 100000)))


def apply_fill_stroke(shape, node, inherited_attrs=None, parent_opacity=1.0):
    """将 SVG 节点的 fill / stroke / opacity / fill-opacity 属性映射到 PPT 形状。
    支持 inherited_attrs 回退（来自父级 <g>）。
    v5.0: 新增 fill-opacity / stroke-opacity 支持。
    """
    if inherited_attrs is None:
        inherited_attrs = {}

    fill_val = node.attrib.get('fill')
    if not fill_val:
        fill_val = inherited_attrs.get('fill')

    stroke_val = node.attrib.get('stroke')
    if not stroke_val:
        stroke_val = inherited_attrs.get('stroke')

    stroke_width = node.attrib.get('stroke-width')
    if not stroke_width:
        stroke_width = inherited_attrs.get('stroke-width', '1')

    stroke_dasharray = node.attrib.get('stroke-dasharray')
    if not stroke_dasharray and inherited_attrs:
        stroke_dasharray = inherited_attrs.get('stroke-dasharray')

    # v4.3: opacity 支持
    opacity_val = node.attrib.get('opacity')
    if not opacity_val:
        opacity_val = inherited_attrs.get('opacity')

    # v5.0: fill-opacity 支持
    fill_opacity_val = node.attrib.get('fill-opacity')
    if not fill_opacity_val:
        fill_opacity_val = inherited_attrs.get('fill-opacity')

    # v5.0: stroke-opacity 支持
    stroke_opacity_val = node.attrib.get('stroke-opacity')
    if not stroke_opacity_val:
        stroke_opacity_val = inherited_attrs.get('stroke-opacity')

    if hasattr(shape, 'fill'):
        if fill_val and fill_val != 'none':
            shape.fill.solid()
            c = parse_color(fill_val)
            if c:
                shape.fill.fore_color.rgb = c
            # 计算最终 fill alpha
            final_alpha = 1.0 * parent_opacity
            if opacity_val:
                try:
                    final_alpha *= float(opacity_val)
                except:
                    pass
            if fill_opacity_val:
                try:
                    final_alpha *= float(fill_opacity_val)
                except:
                    pass
            if final_alpha < 1.0:
                spPr = _find_spPr(shape)
                _apply_alpha(spPr, final_alpha)
        else:
            try:
                spPr = _find_spPr(shape)
                if spPr is not None:
                    for old_fill in list(spPr):
                        tl = old_fill.tag.split('}')[-1] if '}' in old_fill.tag else old_fill.tag
                        if 'Fill' in tl:
                            spPr.remove(old_fill)
                    etree.SubElement(spPr, qn('a:noFill'))
            except:
                pass

    if hasattr(shape, 'line'):
        if stroke_val and stroke_val != 'none':
            c = parse_color(stroke_val)
            if c:
                shape.line.color.rgb = c
            sw = float(stroke_width)
            shape.line.width = Pt(sw * FONT_SCALE * 0.6)
            if stroke_dasharray:
                # Kept for legacy SVGs, but new generated SVGs should avoid it.
                warn("stroke-dasharray encountered; converted to PPT dash style, but current SVG rules prohibit dasharray")
                shape.line.dash_style = MSO_LINE.DASH

            so = 1.0 * parent_opacity
            if stroke_opacity_val:
                try:
                    so *= float(stroke_opacity_val)
                except:
                    pass
            if so < 1.0:
                try:
                    ln_elem = shape._element.find('.//' + qn('a:ln'))
                    if ln_elem is None:
                        sp_el = _find_spPr(shape)
                        if sp_el is not None:
                            ln_elem = sp_el.find(qn('a:ln'))
                    if ln_elem is not None:
                        sf = ln_elem.find(qn('a:solidFill'))
                        if sf is not None:
                            clr = sf.find(qn('a:srgbClr'))
                            if clr is not None:
                                for old_a in clr.findall(qn('a:alpha')):
                                    clr.remove(old_a)
                                alpha_el = etree.SubElement(clr, qn('a:alpha'))
                                alpha_el.set('val', str(int(so * 100000)))
                except:
                    pass
        else:
            try:
                spPr = _find_spPr(shape)
                if spPr is not None:
                    ln_elem = spPr.find(qn('a:ln'))
                    if ln_elem is None:
                        ln_elem = etree.SubElement(spPr, qn('a:ln'))
                    else:
                        for old_fill in list(ln_elem):
                            if 'Fill' in old_fill.tag:
                                ln_elem.remove(old_fill)
                    etree.SubElement(ln_elem, qn('a:noFill'))
            except:
                pass

    try:
        remove_shadow(shape)
    except:
        pass




def svg_to_inches(v):
    return v * SCALE


def estimate_text_width(text, font_size_pt):
    """粗略估算文本宽度 (inches)。"""
    if not text:
        return 0
    w = 0
    for ch in text:
        if ord(ch) > 0x2E7F:
            w += font_size_pt
        else:
            w += font_size_pt * 0.58
    return w / 72.0


def parse_clip_paths(root):
    """解析 SVG <defs> 中的 <clipPath>，提取 clipPath id → rx 映射。"""
    clip_rx_map = {}
    for node in root.iter():
        tag = node.tag.split('}')[-1] if '}' in node.tag else node.tag
        if tag == 'clipPath':
            cp_id = node.attrib.get('id')
            if cp_id:
                for child in node:
                    ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if ctag == 'rect':
                        rx = child.attrib.get('rx')
                        if rx:
                            clip_rx_map[cp_id] = rx
    return clip_rx_map


def parse_css_inset_clip_path(clip_path_attr):
    """v4.3: 解析 CSS clip-path="inset(... round R1 R2 R3 R4)" 语法。
    返回 max(R1..R4) 作为圆角半径，没有 round 则返回 None。
    """
    if not clip_path_attr or 'inset' not in clip_path_attr:
        return None
    m = re.search(r'round\s+([\d.]+)(?:\s+[\d.]+)*', clip_path_attr)
    if m:
        round_part = clip_path_attr[m.start():].replace(')', '')
        nums = re.findall(r'[\d.]+', round_part)
        if nums:
            return str(max(float(n) for n in nums))
    return None


def set_slide_background(slide, color_rgb):
    """设置 slide 背景为纯色。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb


def add_missing_image_placeholder(slide, x, y, w, h, href):
    """Draw a visible placeholder when legacy/test conversion allows missing images."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(204, 204, 204)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Missing image\n{href or '(no href)'}"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(153, 153, 153)


# ═══════════════════════════════════════
# SVG <path> 解析器 (v5.0)
# ═══════════════════════════════════════

def _tokenize_path(d):
    """将 SVG path d 属性拆分为 token 列表。"""
    # 分割命令字母和数字（含负号和小数点）
    tokens = re.findall(r'[MmLlHhVvCcSsQqTtAaZz]|[-+]?[0-9]*\.?[0-9]+(?:[eE][-+]?[0-9]+)?', d)
    return tokens


def _parse_path_to_points(d, offset_x=0, offset_y=0, scale_x=1.0, scale_y=1.0):
    """解析 SVG path d 属性为一系列顶点坐标（绝对坐标）。
    支持：M, L, H, V, C, Q, S, T, A, Z（大写=绝对，小写=相对）。
    曲线命令用折线近似。
    返回：列表的列表（每个子路径一个列表）。
    """
    tokens = _tokenize_path(d)
    subpaths = []  # 所有子路径
    current_points = []  # 当前子路径的点
    cx, cy = 0.0, 0.0  # 当前点
    start_x, start_y = 0.0, 0.0  # 子路径起点
    i = 0
    last_cmd = ''
    last_control_x, last_control_y = 0.0, 0.0  # 上一个控制点（用于 S/T）

    def num():
        nonlocal i
        if i < len(tokens):
            val = float(tokens[i])
            i += 1
            return val
        return 0.0

    def add_point(x, y):
        nonlocal cx, cy
        ax = (x + offset_x) * scale_x
        ay = (y + offset_y) * scale_y
        current_points.append((ax, ay))
        cx, cy = x, y

    def bezier_cubic(x0, y0, x1, y1, x2, y2, x3, y3, steps=8):
        """三次贝塞尔曲线折线近似。"""
        pts = []
        for s in range(1, steps + 1):
            t = s / steps
            t2 = t * t
            t3 = t2 * t
            mt = 1 - t
            mt2 = mt * mt
            mt3 = mt2 * mt
            px = mt3 * x0 + 3 * mt2 * t * x1 + 3 * mt * t2 * x2 + t3 * x3
            py = mt3 * y0 + 3 * mt2 * t * y1 + 3 * mt * t2 * y2 + t3 * y3
            pts.append((px, py))
        return pts

    def bezier_quad(x0, y0, x1, y1, x2, y2, steps=6):
        """二次贝塞尔曲线折线近似。"""
        pts = []
        for s in range(1, steps + 1):
            t = s / steps
            mt = 1 - t
            px = mt * mt * x0 + 2 * mt * t * x1 + t * t * x2
            py = mt * mt * y0 + 2 * mt * t * y1 + t * t * y2
            pts.append((px, py))
        return pts

    def arc_to_points(rx_a, ry_a, x_rot, large_arc, sweep, ex, ey, steps=12):
        """椭圆弧折线近似（简化版：用直线连接起终点的弧段）。"""
        # 简化实现：用二次贝塞尔近似弧
        mx = (cx + ex) / 2
        my = (cy + ey) / 2
        # 粗略控制点偏移
        dx = ex - cx
        dy = ey - cy
        dist = math.sqrt(dx*dx + dy*dy)
        if dist < 0.01:
            return [(ex, ey)]
        # 使用弧的半径来估算弯曲程度
        bulge = min(rx_a, ry_a) * 0.5
        if sweep:
            bulge = -bulge
        # 法线方向
        nx = -dy / dist * bulge
        ny = dx / dist * bulge
        ctrl_x = mx + nx
        ctrl_y = my + ny
        return bezier_quad(cx, cy, ctrl_x, ctrl_y, ex, ey, steps)

    while i < len(tokens):
        token = tokens[i]

        if token in 'MmLlHhVvCcSsQqTtAaZz':
            cmd = token
            i += 1
        else:
            # 隐式重复上一个命令
            cmd = last_cmd
            if not cmd:
                i += 1
                continue

        if cmd == 'M':
            if current_points:
                subpaths.append(current_points)
                current_points = []
            x, y = num(), num()
            start_x, start_y = x, y
            add_point(x, y)
            last_cmd = 'L'  # 后续隐式为 L
            last_control_x, last_control_y = x, y
        elif cmd == 'm':
            if current_points:
                subpaths.append(current_points)
                current_points = []
            dx, dy = num(), num()
            x, y = cx + dx, cy + dy
            start_x, start_y = x, y
            add_point(x, y)
            last_cmd = 'l'
            last_control_x, last_control_y = x, y
        elif cmd == 'L':
            x, y = num(), num()
            add_point(x, y)
            last_cmd = 'L'
            last_control_x, last_control_y = x, y
        elif cmd == 'l':
            dx, dy = num(), num()
            add_point(cx + dx, cy + dy)
            last_cmd = 'l'
            last_control_x, last_control_y = cx, cy
        elif cmd == 'H':
            x = num()
            add_point(x, cy)
            last_cmd = 'H'
            last_control_x, last_control_y = x, cy
        elif cmd == 'h':
            dx = num()
            add_point(cx + dx, cy)
            last_cmd = 'h'
            last_control_x, last_control_y = cx, cy
        elif cmd == 'V':
            y = num()
            add_point(cx, y)
            last_cmd = 'V'
            last_control_x, last_control_y = cx, y
        elif cmd == 'v':
            dy = num()
            add_point(cx, cy + dy)
            last_cmd = 'v'
            last_control_x, last_control_y = cx, cy
        elif cmd == 'C':
            x1, y1 = num(), num()
            x2, y2 = num(), num()
            x3, y3 = num(), num()
            pts = bezier_cubic(cx, cy, x1, y1, x2, y2, x3, y3)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'C'
            last_control_x, last_control_y = x2, y2
        elif cmd == 'c':
            dx1, dy1 = num(), num()
            dx2, dy2 = num(), num()
            dx3, dy3 = num(), num()
            x1, y1 = cx + dx1, cy + dy1
            x2, y2 = cx + dx2, cy + dy2
            x3, y3 = cx + dx3, cy + dy3
            pts = bezier_cubic(cx, cy, x1, y1, x2, y2, x3, y3)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'c'
            last_control_x, last_control_y = x2, y2
        elif cmd == 'S':
            # 反射控制点
            rx = 2 * cx - last_control_x
            ry = 2 * cy - last_control_y
            x2, y2 = num(), num()
            x3, y3 = num(), num()
            pts = bezier_cubic(cx, cy, rx, ry, x2, y2, x3, y3)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'S'
            last_control_x, last_control_y = x2, y2
        elif cmd == 's':
            rx = 2 * cx - last_control_x
            ry = 2 * cy - last_control_y
            dx2, dy2 = num(), num()
            dx3, dy3 = num(), num()
            x2, y2 = cx + dx2, cy + dy2
            x3, y3 = cx + dx3, cy + dy3
            pts = bezier_cubic(cx, cy, rx, ry, x2, y2, x3, y3)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 's'
            last_control_x, last_control_y = x2, y2
        elif cmd == 'Q':
            x1, y1 = num(), num()
            x2, y2 = num(), num()
            pts = bezier_quad(cx, cy, x1, y1, x2, y2)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'Q'
            last_control_x, last_control_y = x1, y1
        elif cmd == 'q':
            dx1, dy1 = num(), num()
            dx2, dy2 = num(), num()
            x1, y1 = cx + dx1, cy + dy1
            x2, y2 = cx + dx2, cy + dy2
            pts = bezier_quad(cx, cy, x1, y1, x2, y2)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'q'
            last_control_x, last_control_y = x1, y1
        elif cmd == 'T':
            rx = 2 * cx - last_control_x
            ry = 2 * cy - last_control_y
            x2, y2 = num(), num()
            pts = bezier_quad(cx, cy, rx, ry, x2, y2)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'T'
            last_control_x, last_control_y = rx, ry
        elif cmd == 't':
            rx = 2 * cx - last_control_x
            ry = 2 * cy - last_control_y
            dx2, dy2 = num(), num()
            x2, y2 = cx + dx2, cy + dy2
            pts = bezier_quad(cx, cy, rx, ry, x2, y2)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 't'
            last_control_x, last_control_y = rx, ry
        elif cmd == 'A':
            arx, ary = num(), num()
            x_rot = num()
            large_arc = int(num())
            sweep = int(num())
            ex, ey = num(), num()
            pts = arc_to_points(arx, ary, x_rot, large_arc, sweep, ex, ey)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'A'
            last_control_x, last_control_y = ex, ey
        elif cmd == 'a':
            arx, ary = num(), num()
            x_rot = num()
            large_arc = int(num())
            sweep = int(num())
            dex, dey = num(), num()
            ex, ey = cx + dex, cy + dey
            pts = arc_to_points(arx, ary, x_rot, large_arc, sweep, ex, ey)
            for px, py in pts:
                add_point(px, py)
            last_cmd = 'a'
            last_control_x, last_control_y = ex, ey
        elif cmd in ('Z', 'z'):
            # 闭合：回到起点
            if current_points:
                add_point(start_x, start_y)
            last_cmd = cmd
            last_control_x, last_control_y = start_x, start_y
        else:
            i += 1

    if current_points:
        subpaths.append(current_points)

    return subpaths


def _path_has_close_command(d):
    """Return True when a path includes an explicit close command."""
    return bool(re.search(r'[Zz]', d or ''))


def _node_fill_value(node, inherited_attrs=None):
    if inherited_attrs is None:
        inherited_attrs = {}
    return node.attrib.get('fill') or inherited_attrs.get('fill')


# ═══════════════════════════════════════
# SVG 节点递归解析 (v5.0)
# ═══════════════════════════════════════

def add_elements(slide, parent_node, offset_x=0, offset_y=0,
                 parent_opacity=1.0, inherited_attrs=None,
                 clip_rx_map=None, scale_x=1.0, scale_y=1.0):
    """递归解析 SVG 节点。v5.0: 新增 scale_x/scale_y 参数。"""
    if inherited_attrs is None:
        inherited_attrs = {}
    if clip_rx_map is None:
        clip_rx_map = {}

    for node in parent_node:
        tag = node.tag.split('}')[-1]

        if tag in ('defs', 'animate', 'animateTransform', 'animateMotion', 'set'):
            continue

        if tag == 'g':
            transform = node.attrib.get('transform', '')
            dx, dy = 0, 0
            sx, sy = 1.0, 1.0
            unsupported = []
            for name in ('rotate', 'skewX', 'skewY', 'matrix'):
                if f'{name}(' in transform:
                    unsupported.append(name)
            if unsupported:
                warn(f"unsupported transform ignored on <g>: {', '.join(unsupported)} in {transform!r}")
            if 'translate' in transform:
                try:
                    m = re.search(r'translate\(\s*([-\d.]+)[\s,]+([-\d.]+)\s*\)', transform)
                    if m:
                        dx = float(m.group(1))
                        dy = float(m.group(2))
                    else:
                        m = re.search(r'translate\(\s*([-\d.]+)\s*\)', transform)
                        if m:
                            dx = float(m.group(1))
                except:
                    pass
            # v5.0: 解析 scale()
            if 'scale' in transform:
                try:
                    m = re.search(r'scale\(\s*([-\d.]+)(?:[\s,]+([-\d.]+))?\s*\)', transform)
                    if m:
                        sx = float(m.group(1))
                        sy = float(m.group(2)) if m.group(2) else sx
                except:
                    pass

            g_opacity = float(node.attrib.get('opacity', 1.0))

            new_inherited = dict(inherited_attrs)
            for attr_name in ('fill', 'font-size', 'font-weight', 'stroke', 'stroke-width',
                              'fill-opacity', 'stroke-opacity'):
                val = node.attrib.get(attr_name)
                if val:
                    new_inherited[attr_name] = val

            # 应用 scale 到 dx/dy（scale 影响平移后的位置）
            new_ox = offset_x + dx
            new_oy = offset_y + dy
            new_sx = scale_x * sx
            new_sy = scale_y * sy

            add_elements(slide, node, new_ox, new_oy,
                         parent_opacity * g_opacity, new_inherited,
                         clip_rx_map, new_sx, new_sy)

        elif tag == 'rect':
            raw_x = float(node.attrib.get('x', 0)) + offset_x
            raw_y = float(node.attrib.get('y', 0)) + offset_y
            raw_w = float(node.attrib.get('width', 0))
            raw_h = float(node.attrib.get('height', 0))

            # 应用 scale
            raw_x *= scale_x
            raw_y *= scale_y
            raw_w *= scale_x
            raw_h *= scale_y

            # ── v4.2: 全画布背景 rect → 设为 slide background ──
            # Only opaque full-canvas rects are true slide backgrounds. A later
            # translucent full-canvas rect is often a texture/wash overlay; if
            # we treat it as the background it overwrites the real base color.
            if raw_w >= (SVG_W - 20) and raw_h >= (SVG_H - 20) and raw_x <= 10 and raw_y <= 10:
                fill_val = node.attrib.get('fill')
                if not fill_val:
                    fill_val = inherited_attrs.get('fill')
                rect_opacity = parent_opacity
                try:
                    rect_opacity *= float(node.attrib.get('opacity', 1.0))
                except:
                    pass
                try:
                    rect_opacity *= float(node.attrib.get('fill-opacity', inherited_attrs.get('fill-opacity', 1.0)))
                except:
                    pass
                if fill_val and fill_val != 'none' and rect_opacity >= 0.99:
                    c = parse_color(fill_val)
                    if c:
                        set_slide_background(slide, c)
                    continue

            x = svg_to_inches(raw_x)
            y = svg_to_inches(raw_y)
            w = svg_to_inches(raw_w)
            h = svg_to_inches(raw_h)
            if w <= 0.01 or h <= 0.01:
                continue

            # ── v4.3: 检查 clip-path 引用的 rx ──
            rx = node.attrib.get('rx')
            if not rx:
                clip_path = node.attrib.get('clip-path', '')
                m = re.match(r'url\(#(.+?)\)', clip_path)
                if m:
                    cp_id = m.group(1)
                    rx = clip_rx_map.get(cp_id)
                if not rx:
                    rx = parse_css_inset_clip_path(clip_path)

            mso = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
            shape = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
            apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)
            if rx:
                try:
                    r_val = svg_to_inches(float(rx) * min(scale_x, scale_y))
                    adj = min(r_val / min(w, h), 0.5)
                    shape.adjustments[0] = adj
                except:
                    pass

        elif tag == 'circle':
            cx = svg_to_inches((float(node.attrib.get('cx', 0)) + offset_x) * scale_x)
            cy = svg_to_inches((float(node.attrib.get('cy', 0)) + offset_y) * scale_y)
            r = svg_to_inches(float(node.attrib.get('r', 0)) * min(scale_x, scale_y))
            if r <= 0.01:
                continue
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - r), Inches(cy - r),
                Inches(r * 2), Inches(r * 2)
            )
            apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)

        elif tag == 'ellipse':
            # v5.0: 椭圆支持
            ecx = svg_to_inches((float(node.attrib.get('cx', 0)) + offset_x) * scale_x)
            ecy = svg_to_inches((float(node.attrib.get('cy', 0)) + offset_y) * scale_y)
            erx = svg_to_inches(float(node.attrib.get('rx', 0)) * scale_x)
            ery = svg_to_inches(float(node.attrib.get('ry', 0)) * scale_y)
            if erx <= 0.01 or ery <= 0.01:
                continue
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(ecx - erx), Inches(ecy - ery),
                Inches(erx * 2), Inches(ery * 2)
            )
            apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)

        elif tag == 'line':
            x1 = svg_to_inches((float(node.attrib.get('x1', 0)) + offset_x) * scale_x)
            y1 = svg_to_inches((float(node.attrib.get('y1', 0)) + offset_y) * scale_y)
            x2 = svg_to_inches((float(node.attrib.get('x2', 0)) + offset_x) * scale_x)
            y2 = svg_to_inches((float(node.attrib.get('y2', 0)) + offset_y) * scale_y)
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x1), Inches(y1), Inches(x2), Inches(y2)
            )
            apply_fill_stroke(connector, node, inherited_attrs, parent_opacity)
            try:
                remove_shadow(connector)
            except:
                pass

        elif tag == 'polygon':
            pts = node.attrib.get('points', '').replace(',', ' ').split()
            coords = []
            for i in range(0, len(pts), 2):
                if i + 1 < len(pts):
                    px = svg_to_inches((float(pts[i]) + offset_x) * scale_x)
                    py = svg_to_inches((float(pts[i + 1]) + offset_y) * scale_y)
                    coords.append((Inches(px), Inches(py)))
            if len(coords) >= 3:
                builder = slide.shapes.build_freeform(coords[0][0], coords[0][1])
                builder.add_line_segments(coords[1:] + [coords[0]])
                shape = builder.convert_to_shape()

                apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)

        elif tag == 'path':
            # v5.0: <path> 元素支持
            d = node.attrib.get('d', '')
            if not d:
                continue

            subpaths = _parse_path_to_points(d, offset_x, offset_y, scale_x, scale_y)
            fill_val = _node_fill_value(node, inherited_attrs)
            path_is_closed = _path_has_close_command(d)
            should_keep_close = path_is_closed and fill_val and fill_val != 'none'

            for pts in subpaths:
                if len(pts) < 2:
                    continue
                # 转换为 Inches
                inch_pts = [(Inches(svg_to_inches(px)), Inches(svg_to_inches(py)))
                            for px, py in pts]
                builder = slide.shapes.build_freeform(inch_pts[0][0], inch_pts[0][1])
                builder.add_line_segments(inch_pts[1:])
                shape = builder.convert_to_shape()
                # Open paths must stay open; filled closed paths need close preserved
                # or PowerPoint can lose/alter the fill.
                if not should_keep_close:
                    for close_el in shape._element.findall('.//' + qn('a:close')):
                        close_el.getparent().remove(close_el)
                apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)

        elif tag == 'image':
            raw_x = (float(node.attrib.get('x', 0)) + offset_x) * scale_x
            raw_y = (float(node.attrib.get('y', 0)) + offset_y) * scale_y
            raw_w = float(node.attrib.get('width', 0)) * scale_x
            raw_h = float(node.attrib.get('height', 0)) * scale_y
            if raw_w <= 0.01 or raw_h <= 0.01:
                continue

            href = node.attrib.get('href') or node.attrib.get('{http://www.w3.org/1999/xlink}href')
            img_path = resolve_image_href(href)
            x = svg_to_inches(raw_x)
            y = svg_to_inches(raw_y)
            w = svg_to_inches(raw_w)
            h = svg_to_inches(raw_h)
            if not img_path or not os.path.isfile(img_path):
                msg = f'image not found: {href}'
                if _strict_missing_images:
                    error(msg)
                else:
                    warn(msg + " — inserted placeholder")
                    add_missing_image_placeholder(slide, x, y, w, h, href)
                continue

            pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            try:
                remove_shadow(pic)
            except:
                pass

        elif tag == 'text':
            _add_text_element(slide, node, offset_x, offset_y, inherited_attrs,
                              scale_x, scale_y)

        elif tag == 'polyline':
            pts = node.attrib.get('points', '').replace(',', ' ').split()
            coords = []
            for i in range(0, len(pts), 2):
                if i + 1 < len(pts):
                    px = svg_to_inches((float(pts[i]) + offset_x) * scale_x)
                    py = svg_to_inches((float(pts[i + 1]) + offset_y) * scale_y)
                    coords.append((Inches(px), Inches(py)))
            if len(coords) >= 2:
                builder = slide.shapes.build_freeform(coords[0][0], coords[0][1])
                builder.add_line_segments(coords[1:])
                shape = builder.convert_to_shape()
                # 强制移除 close，保留开放路径
                for close_el in shape._element.findall('.//' + qn('a:close')):
                    close_el.getparent().remove(close_el)
                apply_fill_stroke(shape, node, inherited_attrs, parent_opacity)


def _apply_letter_spacing(run, spacing_val):
    """v4.3: 通过 XML 注入 a:rPr spc 属性设置 letter-spacing。
    spacing_val: SVG px 值 → PPT 百分之一点 (1/100 pt)
    """
    if not spacing_val:
        return
    try:
        spc = int(float(spacing_val) * 100)  # px → 百分之一点
        rPr = run._r.find(qn('a:rPr'))
        if rPr is None:
            rPr = etree.SubElement(run._r, qn('a:rPr'))
            # 移到第一个位置
            run._r.insert(0, rPr)
        rPr.set('spc', str(spc))
    except:
        pass


def _add_text_element(slide, node, offset_x, offset_y, inherited_attrs=None,
                      scale_x=1.0, scale_y=1.0):
    """处理 <text> 节点，生成精确定位的原生文本框。
    v5.0: 新增 scale_x/scale_y 支持。
    """
    if inherited_attrs is None:
        inherited_attrs = {}

    x_svg = (float(node.attrib.get('x', 0)) + offset_x) * scale_x
    y_svg = (float(node.attrib.get('y', 0)) + offset_y) * scale_y

    fs_svg = float(node.attrib.get('font-size',
                   inherited_attrs.get('font-size', '20')))
    # v5.0: scale 影响字号
    fs_svg *= min(scale_x, scale_y)

    anchor = node.attrib.get('text-anchor', 'start')

    fill_str = node.attrib.get('fill',
               inherited_attrs.get('fill', '#333333'))
    fill = parse_color(fill_str)

    font_weight = node.attrib.get('font-weight',
                  inherited_attrs.get('font-weight', 'normal'))

    # v4.3: letter-spacing
    letter_spacing = node.attrib.get('letter-spacing',
                     inherited_attrs.get('letter-spacing'))

    fs_pt = fs_svg * FONT_SCALE

    full_text = node.text or ''
    for child in node:
        ctag = child.tag.split('}')[-1]
        if ctag == 'tspan':
            full_text += (child.text or '')
        if child.tail:
            full_text += child.tail
    full_text = full_text.strip()
    if not full_text:
        return

    # v4.3: letter-spacing 影响文本宽度估算
    ls_extra = 0
    if letter_spacing:
        try:
            ls_extra = float(letter_spacing) * len(full_text) * SCALE
        except:
            pass
    text_w = estimate_text_width(full_text, fs_pt) * 1.15 + ls_extra
    text_w = max(text_w, 0.5)
    text_h = fs_pt / 72.0 * 1.6

    x_in = svg_to_inches(x_svg)
    y_in = svg_to_inches(y_svg)
    y_top = y_in - (fs_pt / 72.0) * 0.85

    if anchor == 'middle':
        tx = x_in - text_w / 2
        align = PP_ALIGN.CENTER
    elif anchor == 'end':
        tx = x_in - text_w
        align = PP_ALIGN.RIGHT
    else:
        tx = x_in
        align = PP_ALIGN.LEFT

    tx = max(tx, 0)
    y_top = max(y_top, 0)

    tb = slide.shapes.add_textbox(Inches(tx), Inches(y_top), Inches(text_w), Inches(text_h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align

    is_bold = font_weight in ('bold', '900', '800', '700')

    if node.text and node.text.strip():
        run = p.add_run()
        run.text = node.text if node.text.strip() else ''
        run.font.size = Pt(fs_pt)
        run.font.bold = is_bold
        run.font.name = 'Microsoft YaHei'
        if fill:
            run.font.color.rgb = fill
        _apply_letter_spacing(run, letter_spacing)

    for child in node:
        ctag = child.tag.split('}')[-1]
        if ctag == 'tspan':
            if child.text:
                run = p.add_run()
                run.text = child.text
                child_fs = float(child.attrib.get('font-size', fs_svg / min(scale_x, scale_y)))
                child_fs *= min(scale_x, scale_y)
                run.font.size = Pt(child_fs * FONT_SCALE)
                c_weight = child.attrib.get('font-weight', font_weight)
                run.font.bold = c_weight in ('bold', '900', '800', '700')
                c_color = parse_color(child.attrib.get('fill')) or fill
                if c_color:
                    run.font.color.rgb = c_color
                run.font.name = 'Microsoft YaHei'
                child_ls = child.attrib.get('letter-spacing', letter_spacing)
                _apply_letter_spacing(run, child_ls)
        if child.tail and child.tail.strip():
            run = p.add_run()
            run.text = child.tail
            run.font.size = Pt(fs_pt)
            run.font.bold = is_bold
            if fill:
                run.font.color.rgb = fill
            run.font.name = 'Microsoft YaHei'
            _apply_letter_spacing(run, letter_spacing)

    try:
        remove_shadow(tb)
    except:
        pass


# ═══════════════════════════════════════
# 演讲者备注写入
# ═══════════════════════════════════════

def add_speaker_notes(slide, notes_text):
    """向 slide 写入演讲者备注。"""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ═══════════════════════════════════════
# 主函数
# ═══════════════════════════════════════

def find_layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return None


def delete_all_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def main():
    global _gradient_map
    global _current_svg_dir
    global _conversion_errors
    global _conversion_warnings
    global _strict_missing_images

    if os.environ.get("SMART_SVG_EXPORT_APPROVED_BY_PPTFLOW") != "1":
        print(
            "ERROR: direct PPT export is blocked. Run `python scripts/pptflow.py <project_dir> export` after export-ready passes.",
            file=sys.stderr,
        )
        sys.exit(2)

    parser = argparse.ArgumentParser(
        description='Smart SVG -> PPTX Converter v5.0')
    parser.add_argument('svgs', nargs='+', help='SVG files (in slide order)')
    parser.add_argument('-o', '--output', default='final_deck.pptx',
                        help='Output PPTX path')
    parser.add_argument('--auto-size', action='store_true',
                        help='Auto-detect SVG canvas size from first SVG viewBox/width/height')
    parser.add_argument('--match-aspect', action='store_true',
                        help='When auto-size, set slide height to match SVG aspect ratio (based on slide width)')
    parser.add_argument('--slide-width-in', type=float, default=SLIDE_W_IN,
                        help='Optional: slide width (inches)')
    parser.add_argument('--slide-height-in', type=float, default=None,
                        help='Optional: slide height (inches). If omitted with --match-aspect, height follows SVG ratio.')
    parser.add_argument('--template', default=None,
                        help='Optional: existing .pptx template')
    parser.add_argument('--layout', default=None,
                        help='Optional: layout name')
    parser.add_argument('--notes', default=None,
                        help='Optional: JSON file mapping SVG filenames to speaker notes')
    parser.add_argument('--report', default=None,
                        help='Optional: write JSON conversion report to this path')
    parser.add_argument('--strict-missing-images', action='store_true',
                        help='Fail conversion when an <image> file is missing. By default, missing images become visible placeholders and are reported.')
    args = parser.parse_args()
    _conversion_errors = []
    _conversion_warnings = []
    _strict_missing_images = bool(args.strict_missing_images)

    svg_files = []
    for pattern in args.svgs:
        expanded = glob.glob(pattern)
        svg_files.extend(sorted(expanded) if expanded else [pattern])

    if not svg_files:
        print('ERROR: No SVG files found.')
        sys.exit(1)

    # ── 画布尺寸自动识别（可选） ──
    if args.auto_size:
        svg_w, svg_h = parse_svg_canvas_size(svg_files[0])
        slide_w_in = args.slide_width_in
        if args.slide_height_in is not None:
            slide_h_in = args.slide_height_in
        elif args.match_aspect and svg_w > 0:
            slide_h_in = slide_w_in * (svg_h / svg_w)
        else:
            slide_h_in = SLIDE_H_IN
        set_canvas(svg_w, svg_h, slide_w_in=slide_w_in, slide_h_in=slide_h_in)
        print(f'Auto canvas: SVG {int(svg_w)}x{int(svg_h)} → Slide {SLIDE_W_IN:.3f}x{SLIDE_H_IN:.3f} in')

    notes_map = {}
    if args.notes:
        try:
            with open(args.notes, 'r', encoding='utf-8') as f:
                notes_map = json.load(f)
            print(f'Loaded speaker notes for {len(notes_map)} slide(s)')
        except Exception as e:
            print(f'WARNING: Could not load notes file: {e}')

    print(f'Found {len(svg_files)} SVG file(s):')
    for f in svg_files:
        print(f'  {f}')

    # ── 模板感知 ──
    if args.template and os.path.isfile(args.template):
        print(f'Using template: {args.template}')
        prs = Presentation(args.template)
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        old_count = len(prs.slides)
        delete_all_slides(prs)
        print(f'  Removed {old_count} existing slide(s) from template')
    else:
        if args.template:
            print(f'WARNING: Template file not found: {args.template}, using blank.')
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    # ── 版式选择 ──
    target_layout = None
    if args.layout:
        target_layout = find_layout_by_name(prs, args.layout)
        if target_layout:
            print(f'Using layout: "{args.layout}"')
        else:
            print(f'WARNING: Layout "{args.layout}" not found, falling back to blank')

    if target_layout is None:
        for layout in prs.slide_layouts:
            if layout.name in ('空白', 'BLANK', 'Blank'):
                target_layout = layout
                break
        if target_layout is None:
            target_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]

    page_reports = []
    for f in svg_files:
        slide = prs.slides.add_slide(target_layout)
        before_errors = len(_conversion_errors)
        before_warnings = len(_conversion_warnings)

        try:
            _current_svg_dir = os.path.dirname(os.path.abspath(f))
            tree = ET.parse(f)
            root = tree.getroot()
            # v4.2: 预解析 clipPath
            clip_rx_map = parse_clip_paths(root)
            # v5.0: 预解析渐变
            _gradient_map = parse_gradients(root)
            add_elements(slide, root, clip_rx_map=clip_rx_map)
            print(f'  OK: {f}')
        except Exception as e:
            error(f'{f}: {e}')
            import traceback
            traceback.print_exc()

        basename = os.path.basename(f)
        if basename in notes_map:
            try:
                add_speaker_notes(slide, notes_map[basename])
                print(f'  NOTES: {basename}')
            except Exception as e:
                error(f'notes error for {basename}: {e}')

        page_reports.append({
            "file": f,
            "status": "fail" if len(_conversion_errors) > before_errors else "ok",
            "errors": _conversion_errors[before_errors:],
            "warnings": _conversion_warnings[before_warnings:],
        })

    output = args.output
    report = {
        "output": output,
        "page_count": len(svg_files),
        "status": "fail" if _conversion_errors else "ok",
        "errors": _conversion_errors,
        "warnings": _conversion_warnings,
        "pages": page_reports,
    }

    report_path = args.report
    if not report_path:
        base, _ = os.path.splitext(output)
        report_path = base + "_conversion_report.json"
    try:
        os.makedirs(os.path.dirname(os.path.abspath(report_path)), exist_ok=True)
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(report, f, ensure_ascii=False, indent=2)
        print(f'Conversion report: {report_path}')
    except Exception as e:
        print(f'WARNING: Could not write conversion report: {e}')

    if _conversion_errors:
        print("\nEXPORT FAILED: SVG to PPT conversion encountered blocking errors.", file=sys.stderr)
        for msg in _conversion_errors:
            print(f"  - {msg}", file=sys.stderr)
        sys.exit(1)

    try:
        prs.save(output)
    except PermissionError:
        output = output.replace('.pptx', '_v2.pptx')
        prs.save(output)
    print(f'\nSaved: {output}')


if __name__ == '__main__':
    main()
